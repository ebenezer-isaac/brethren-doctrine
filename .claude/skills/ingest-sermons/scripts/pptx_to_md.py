"""Extract text and speaker notes from a .pptx into a flat markdown file.

Usage: python pptx_to_md.py <input.pptx> <output.md>

Loses visual layout. Flag visually-rich slides in parsing_notes downstream.
Requires: pip install python-pptx
"""
import sys
import os
from pptx import Presentation


def main(src: str, dst: str) -> None:
    prs = Presentation(src)
    out: list[str] = [f"# {os.path.basename(src)}\n"]
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"\n## Slide {i}\n")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        out.append(text)
            elif shape.shape_type == 13:
                out.append(f"[image: {shape.name}]")
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                out.append(f"\n**Speaker notes:** {notes}")
    os.makedirs(os.path.dirname(dst) or ".", exist_ok=True)
    with open(dst, "w", encoding="utf-8") as f:
        f.write("\n".join(out))
    print(f"wrote {dst} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        sys.exit("usage: pptx_to_md.py <input.pptx> <output.md>")
    main(sys.argv[1], sys.argv[2])
